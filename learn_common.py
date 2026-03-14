#!/usr/bin/env python3
"""
Shared utilities for the learner agent system.
Table formatting, text extractors, DB helpers, chunking, LLM extraction.
"""
import io
import json
import os
import re
import shutil
import sqlite3
import sys
import time

import requests
from PyPDF2 import PdfReader

from learner_config import get_env_path, get_db_path, get_model, get_chunk_size

# ── Config ──────────────────────────────────────────────────────────────────
ENV_PATH = get_env_path()
DB_PATH = get_db_path()
MODEL = get_model()
CHUNK_SIZE = get_chunk_size()

# ── Env Loading ─────────────────────────────────────────────────────────────

def load_env():
    env = {}
    with open(ENV_PATH) as f:
        for line in f:
            line = line.strip()
            if line and "=" in line and not line.startswith("#"):
                if line.startswith("export "):
                    line = line[7:]
                k, v = line.split("=", 1)
                env[k.strip()] = v.strip().strip('"').strip("'")
    return env


# ── Table Output (Sticky Headers) ──────────────────────────────────────────

def _col_widths():
    cols = shutil.get_terminal_size().columns
    fixed = 10 + 6 + 8 + 9 + 12  # progress + chunks + insights + size + gaps
    remaining = max(cols - fixed, 40)
    name_w = max(remaining * 30 // 100, 16)
    loc_w = remaining - name_w
    return {"progress": 10, "name": name_w, "location": loc_w, "chunks": 6, "insights": 8, "size": 9}


def format_size(nbytes):
    if nbytes < 1024:
        return f"{nbytes} B"
    elif nbytes < 1024 * 1024:
        return f"{nbytes / 1024:.1f} KB"
    else:
        return f"{nbytes / (1024 * 1024):.1f} MB"


def _tbl_sep():
    w = _col_widths()
    return f"  {'─'*w['progress']}  {'─'*w['name']}  {'─'*w['location']}  {'─'*w['chunks']}  {'─'*w['insights']}  {'─'*w['size']}"


def print_table_header(location_label="Location", sticky=True):
    w = _col_widths()
    header = f"  {'Progress':^{w['progress']}}  {'File Name':<{w['name']}}  {location_label:<{w['location']}}  {'Chunks':^{w['chunks']}}  {'Insights':^{w['insights']}}  {'Size':>{w['size']}}"
    sep = _tbl_sep()
    if sticky and sys.stdout.isatty():
        rows = shutil.get_terminal_size().lines
        sys.stdout.write("\033[1;1H\033[J")   # cursor to top, clear screen
        print(header)
        print(sep)
        sys.stdout.write(f"\033[3;{rows}r")   # scroll region: line 3 to bottom
        sys.stdout.write("\033[3;1H")          # cursor into scroll region
        sys.stdout.flush()
    else:
        print(header)
        print(sep)


def print_table_row(progress, name, location, chunks, insights, size_bytes):
    w = _col_widths()
    name_s = name[:w["name"]-1] + "…" if len(name) > w["name"] else name
    loc_s = location[:w["location"]-1] + "…" if len(location) > w["location"] else location
    size_s = format_size(size_bytes)
    print(f"  {progress:^{w['progress']}}  {name_s:<{w['name']}}  {loc_s:<{w['location']}}  {chunks:^{w['chunks']}}  {insights:^{w['insights']}}  {size_s:>{w['size']}}")


def print_table_footer():
    print(_tbl_sep())
    if sys.stdout.isatty():
        sys.stdout.write("\033[r")  # reset scroll region only, don't jump cursor
        sys.stdout.flush()


# ── Text Extractors ─────────────────────────────────────────────────────────

def extract_text(raw_bytes, ext):
    """Extract text content from raw bytes based on file extension."""
    if ext == "pdf":
        return _extract_pdf(raw_bytes)
    if ext in ("docx", "doc"):
        return _extract_docx(raw_bytes)
    if ext in ("xlsx", "xls"):
        return _extract_xlsx(raw_bytes)
    if ext in ("pptx", "ppt"):
        return _extract_pptx(raw_bytes)
    try:
        return raw_bytes.decode("utf-8", errors="replace")
    except Exception:
        return raw_bytes.decode("latin-1", errors="replace")


def _extract_pdf(raw_bytes):
    reader = PdfReader(io.BytesIO(raw_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages) if pages else None


def _extract_docx(raw_bytes):
    import docx
    doc = docx.Document(io.BytesIO(raw_bytes))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(raw_bytes):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    rows = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
    wb.close()
    return "\n".join(rows) if rows else None


def _extract_pptx(raw_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else None


# ── Database Helpers ────────────────────────────────────────────────────────

def open_db(db_path=None):
    if db_path is None:
        db_path = DB_PATH
    os.makedirs(os.path.dirname(db_path), exist_ok=True)
    conn = sqlite3.connect(db_path)
    conn.execute("PRAGMA journal_mode=WAL")
    return conn


def is_processed(conn, file_path):
    row = conn.execute(
        "SELECT error FROM processed_files WHERE file_path = ?", (file_path,)
    ).fetchone()
    if row is None:
        return False
    if row[0] is None:
        return True
    error = row[0]
    transient_errors = ("LLM", "timeout", "rate limit", "429", "500", "502", "503")
    return not any(t.lower() in error.lower() for t in transient_errors)


def mark_processed(conn, file_path, source, learning_count, error=None, file_size=0):
    conn.execute(
        "INSERT OR REPLACE INTO processed_files (file_path, source, file_size, learning_count, error) "
        "VALUES (?, ?, ?, ?, ?)",
        (file_path, source, file_size, learning_count, error),
    )
    conn.commit()


def mark_error(conn, file_path, source, error_msg):
    conn.execute(
        "INSERT OR REPLACE INTO processed_files (file_path, source, learning_count, error) "
        "VALUES (?, ?, 0, ?)",
        (file_path, source, error_msg),
    )
    conn.commit()


def insert_learnings(conn, learnings_list, source, agent, folder, file_path, file_name):
    for learning in learnings_list:
        conn.execute(
            "INSERT OR IGNORE INTO learnings (source, agent, folder, file_path, file_name, learning) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (source, agent, folder, file_path, file_name, learning),
        )
    conn.commit()


# ── Chunking + LLM Extraction ──────────────────────────────────────────────

def parse_learnings(text):
    if not text:
        return []
    cleaned = text.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    cleaned = cleaned.strip()
    try:
        result = json.loads(cleaned)
        if isinstance(result, list):
            return [s.strip() for s in result if isinstance(s, str) and len(s.strip()) > 10]
    except json.JSONDecodeError:
        pass
    learnings = []
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        for prefix in ("- ", "* ", "• ", "– "):
            if line.startswith(prefix):
                line = line[len(prefix):]
                break
        else:
            m = re.match(r"^\d+[.)]\s*", line)
            if m:
                line = line[m.end():]
            else:
                continue
        line = line.strip()
        if len(line) > 10:
            learnings.append(line)
    return learnings


def chunk_content(text, chunk_size=None):
    if chunk_size is None:
        chunk_size = CHUNK_SIZE
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    paragraphs = text.split("\n\n")
    current = ""
    for para in paragraphs:
        if len(current) + len(para) + 2 > chunk_size and current:
            chunks.append(current)
            current = para
        else:
            current = current + "\n\n" + para if current else para
    if current:
        chunks.append(current)
    return chunks if chunks else [text[:chunk_size]]


def _llm_call(prompt, openrouter_key, retries=3):
    """Single LLM call with retries. Returns parsed learnings list."""
    for attempt in range(retries):
        try:
            r = requests.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {openrouter_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": MODEL,
                    "messages": [{"role": "user", "content": prompt}],
                    "max_tokens": 1500,
                },
                timeout=90,
            )
            if r.status_code == 200:
                txt = r.json().get("choices", [{}])[0].get("message", {}).get("content", "")
                return parse_learnings(txt)
            elif r.status_code == 429:
                wait = 2 ** (attempt + 1)
                print(f"    Rate limited, waiting {wait}s...")
                time.sleep(wait)
                continue
            else:
                print(f"    LLM error {r.status_code}: {r.text[:100]}")
        except Exception as e:
            print(f"    LLM exception (attempt {attempt + 1}): {e}")
            if attempt < retries - 1:
                time.sleep(2 ** (attempt + 1))
    return []


def extract_learnings_chunked(content, prompt_builder, openrouter_key, on_chunk=None):
    """
    Extract learnings with verbose chunk progress.

    Args:
        content: Text content to process.
        prompt_builder: callable(chunk_text) -> full prompt string.
        openrouter_key: API key.
        on_chunk: optional callback(chunk_index, total_chunks, chunk_learnings) -- called per chunk.

    Returns: (all_learnings, num_chunks)
    """
    if not content or len(content) < 20:
        return [], 1

    chunks = chunk_content(content)
    num_chunks = len(chunks)
    if num_chunks > 1:
        print(f"    Large content: splitting into {num_chunks} chunks")

    all_learnings = []
    for i, chunk in enumerate(chunks):
        prompt = prompt_builder(chunk)
        chunk_learnings = _llm_call(prompt, openrouter_key)
        all_learnings.extend(chunk_learnings)
        if on_chunk and num_chunks > 1:
            on_chunk(i + 1, num_chunks, chunk_learnings)

    seen = set()
    return [l for l in all_learnings if not (l in seen or seen.add(l))], num_chunks
